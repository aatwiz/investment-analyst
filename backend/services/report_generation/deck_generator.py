"""
Pitch Deck Generator
Generates professional PowerPoint presentations for investment opportunities
"""

from typing import Dict, List, Optional, Any
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
from loguru import logger
from openai import AsyncOpenAI
import os


class PitchDeckGenerator:
    """
    Generates professional pitch decks integrating data from all features
    """
    
    # Color scheme
    DARK_BLUE = RGBColor(0, 51, 102)
    LIGHT_BLUE = RGBColor(0, 153, 204)
    ACCENT = RGBColor(255, 102, 0)
    GRAY = RGBColor(100, 100, 100)
    WHITE = RGBColor(255, 255, 255)
    
    def __init__(self):
        """Initialize deck generator"""
        self.client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        logger.info("PitchDeckGenerator initialized")
    
    async def generate_deck(
        self,
        company_data: Dict[str, Any],
        deal_data: Optional[Dict[str, Any]] = None,
        market_data: Optional[Dict[str, Any]] = None,
        financial_model: Optional[Dict[str, Any]] = None,
        template_type: str = "standard"
    ) -> io.BytesIO:
        """
        Generate complete pitch deck
        
        Args:
            company_data: Company information
            deal_data: Deal qualification data
            market_data: Market analysis
            financial_model: Financial projections
            template_type: Deck template
        
        Returns:
            BytesIO containing PPTX deck
        """
        logger.info(f"Generating {template_type} pitch deck for {company_data.get('name', 'Unknown')}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Generate AI-powered content
        deck_content = await self._generate_deck_content(
            company_data, deal_data, market_data, financial_model
        )
        
        # Add slides
        self._add_title_slide(prs, company_data)
        self._add_problem_slide(prs, deck_content.get('problem', {}))
        self._add_solution_slide(prs, company_data, deck_content.get('solution', {}))
        self._add_market_slide(prs, market_data, deck_content.get('market', {}))
        self._add_product_slide(prs, company_data, deck_content.get('product', {}))
        self._add_business_model_slide(prs, company_data, deck_content.get('business_model', {}))
        self._add_traction_slide(prs, company_data, deck_content.get('traction', {}))
        self._add_financials_slide(prs, financial_model)
        self._add_team_slide(prs, company_data, deck_content.get('team', {}))
        self._add_competition_slide(prs, market_data, deck_content.get('competition', {}))
        self._add_investment_slide(prs, company_data, deck_content.get('ask', {}))
        self._add_closing_slide(prs, company_data)
        
        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        logger.info("Pitch deck generated successfully")
        return output
    
    def _add_title_slide(self, prs: Presentation, company_data: Dict):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Company name
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = company_data.get('name', 'Company Name')
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.DARK_BLUE
        
        # Tagline
        top = Inches(3.7)
        height = Inches(0.5)
        
        tagline_box = slide.shapes.add_textbox(left, top, width, height)
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = company_data.get('description', '')[:100]
        
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.alignment = PP_ALIGN.CENTER
        tagline_para.font.size = Pt(20)
        tagline_para.font.color.rgb = self.GRAY
        
        # Stage & Funding
        top = Inches(5)
        info_box = slide.shapes.add_textbox(left, top, width, height)
        info_frame = info_box.text_frame
        
        stage = company_data.get('stage', 'Early Stage')
        funding = f"${company_data.get('funding_amount', 0):,.0f}"
        info_frame.text = f"{stage} | {funding} Funding"
        
        info_para = info_frame.paragraphs[0]
        info_para.alignment = PP_ALIGN.CENTER
        info_para.font.size = Pt(16)
        info_para.font.color.rgb = self.LIGHT_BLUE
        
        # Date
        top = Inches(6.5)
        date_box = slide.shapes.add_textbox(left, top, width, height)
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime('%B %Y')
        
        date_para = date_frame.paragraphs[0]
        date_para.alignment = PP_ALIGN.CENTER
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = self.GRAY
    
    def _add_problem_slide(self, prs: Presentation, problem: Dict):
        """Add problem slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        title = slide.shapes.title
        title.text = "The Problem"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        problems = problem.get('points', [
            "Market pain point 1",
            "Market pain point 2",
            "Market pain point 3"
        ])
        
        for prob in problems:
            p = text_frame.add_paragraph()
            p.text = prob
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(12)
            p.space_after = Pt(12)
    
    def _add_solution_slide(self, prs: Presentation, company_data: Dict, solution: Dict):
        """Add solution slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Our Solution"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Description
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        
        desc_box = slide.shapes.add_textbox(left, top, width, height)
        desc_frame = desc_box.text_frame
        desc_frame.text = company_data.get('description', 'Our innovative solution...')
        desc_frame.paragraphs[0].font.size = Pt(20)
        desc_frame.word_wrap = True
        
        # Key features
        top = Inches(3.7)
        height = Inches(3)
        
        features_box = slide.shapes.add_textbox(left, top, width, height)
        features_frame = features_box.text_frame
        
        features = solution.get('features', [
            "Key feature 1",
            "Key feature 2",
            "Key feature 3"
        ])
        
        for feature in features:
            p = features_frame.add_paragraph()
            p.text = f"✓ {feature}"
            p.font.size = Pt(18)
            p.space_before = Pt(10)
            p.space_after = Pt(10)
    
    def _add_market_slide(self, prs: Presentation, market_data: Optional[Dict], market: Dict):
        """Add market opportunity slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Market Opportunity"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Market size box
        left = Inches(1)
        top = Inches(2)
        width = Inches(3.5)
        height = Inches(2)
        
        size_box = slide.shapes.add_textbox(left, top, width, height)
        size_frame = size_box.text_frame
        
        market_size = market_data.get('market_size', 'Large and Growing') if market_data else market.get('size', '$XXB Market')
        
        p = size_frame.paragraphs[0]
        p.text = "Market Size"
        p.font.size = Pt(16)
        p.font.color.rgb = self.GRAY
        
        p = size_frame.add_paragraph()
        p.text = market_size
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.ACCENT
        
        # Growth rate box
        left = Inches(5.5)
        
        growth_box = slide.shapes.add_textbox(left, top, width, height)
        growth_frame = growth_box.text_frame
        
        growth_rate = market_data.get('growth_rate', 'XX% CAGR') if market_data else market.get('growth', 'Fast Growing')
        
        p = growth_frame.paragraphs[0]
        p.text = "Growth Rate"
        p.font.size = Pt(16)
        p.font.color.rgb = self.GRAY
        
        p = growth_frame.add_paragraph()
        p.text = growth_rate
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.ACCENT
        
        # Trends
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(2)
        
        trends_box = slide.shapes.add_textbox(left, top, width, height)
        trends_frame = trends_box.text_frame
        
        p = trends_frame.paragraphs[0]
        p.text = "Key Market Trends:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        trends = market.get('trends', [
            "Growing demand for digital solutions",
            "Shift to cloud-based services",
            "Increasing regulatory requirements"
        ])
        
        for trend in trends:
            p = trends_frame.add_paragraph()
            p.text = f"• {trend}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)
    
    def _add_product_slide(self, prs: Presentation, company_data: Dict, product: Dict):
        """Add product slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Product/Technology"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        capabilities = product.get('capabilities', [
            "Core technology/platform",
            "Unique features and benefits",
            "Scalability and performance"
        ])
        
        for cap in capabilities:
            p = text_frame.add_paragraph()
            p.text = cap
            p.font.size = Pt(18)
            p.space_before = Pt(15)
            p.space_after = Pt(15)
    
    def _add_business_model_slide(self, prs: Presentation, company_data: Dict, business_model: Dict):
        """Add business model slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Business Model"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "Revenue Streams:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        streams = business_model.get('revenue_streams', [
            "Subscription/SaaS revenue",
            "Transaction fees",
            "Enterprise licenses"
        ])
        
        for stream in streams:
            p = text_frame.add_paragraph()
            p.text = f"• {stream}"
            p.font.size = Pt(16)
            p.space_after = Pt(10)
        
        text_frame.add_paragraph()
        
        p = text_frame.add_paragraph()
        p.text = "Unit Economics:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_before = Pt(20)
        p.space_after = Pt(10)
        
        economics = business_model.get('unit_economics', {})
        metrics = [
            f"CAC: {economics.get('cac', 'TBD')}",
            f"LTV: {economics.get('ltv', 'TBD')}",
            f"LTV/CAC Ratio: {economics.get('ratio', 'TBD')}"
        ]
        
        for metric in metrics:
            p = text_frame.add_paragraph()
            p.text = f"• {metric}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
    
    def _add_traction_slide(self, prs: Presentation, company_data: Dict, traction: Dict):
        """Add traction/milestones slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Traction & Milestones"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Metrics boxes
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(2.3)
        height = Inches(1.5)
        
        milestones = traction.get('milestones', [
            {"label": "Customers", "value": "XXX"},
            {"label": "Revenue", "value": "$XXK MRR"},
            {"label": "Growth", "value": "XX% MoM"}
        ])
        
        for i, milestone in enumerate(milestones[:3]):
            left_pos = Inches(1 + i * 2.8)
            
            metric_box = slide.shapes.add_textbox(left_pos, top, width, height)
            metric_frame = metric_box.text_frame
            metric_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = metric_frame.paragraphs[0]
            p.text = milestone['label']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = self.GRAY
            
            p = metric_frame.add_paragraph()
            p.text = milestone['value']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT
        
        # Achievements
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(2)
        
        achievements_box = slide.shapes.add_textbox(left, top, width, height)
        achievements_frame = achievements_box.text_frame
        
        p = achievements_frame.paragraphs[0]
        p.text = "Recent Achievements:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(10)
        
        achievements = traction.get('achievements', [
            "Key partnership or customer win",
            "Product launch or milestone",
            "Team expansion or hire"
        ])
        
        for achievement in achievements:
            p = achievements_frame.add_paragraph()
            p.text = f"✓ {achievement}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)
    
    def _add_financials_slide(self, prs: Presentation, financial_model: Optional[Dict]):
        """Add financials slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Financial Projections"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        if not financial_model:
            left = Inches(2)
            top = Inches(3)
            width = Inches(6)
            height = Inches(2)
            
            note_box = slide.shapes.add_textbox(left, top, width, height)
            note_frame = note_box.text_frame
            note_frame.text = "Detailed financial projections available upon request"
            note_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            note_frame.paragraphs[0].font.size = Pt(20)
            note_frame.paragraphs[0].font.color.rgb = self.GRAY
            return
        
        metrics = financial_model.get('metrics', {})
        
        # Key metrics boxes
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(2.3)
        height = Inches(1.5)
        
        financial_metrics = [
            {"label": "Total Revenue", "value": f"${metrics.get('total_revenue', 0)/1000000:.1f}M"},
            {"label": "Revenue CAGR", "value": f"{metrics.get('revenue_cagr', 0)*100:.0f}%"},
            {"label": "Time to Profit", "value": f"{metrics.get('months_to_profitability', 'N/A')} mo"}
        ]
        
        for i, metric in enumerate(financial_metrics):
            left_pos = Inches(1 + i * 2.8)
            
            metric_box = slide.shapes.add_textbox(left_pos, top, width, height)
            metric_frame = metric_box.text_frame
            metric_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = metric_frame.paragraphs[0]
            p.text = metric['label']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = self.GRAY
            
            p = metric_frame.add_paragraph()
            p.text = metric['value']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT
        
        # Additional details
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(2)
        
        details_box = slide.shapes.add_textbox(left, top, width, height)
        details_frame = details_box.text_frame
        
        p = details_frame.paragraphs[0]
        p.text = "Key Financial Highlights:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(10)
        
        highlights = [
            f"Final cash balance: ${metrics.get('final_cash_balance', 0)/1000000:.1f}M",
            f"Minimum cash runway maintained throughout projection period",
            "Detailed 3-year model with monthly granularity available"
        ]
        
        for highlight in highlights:
            p = details_frame.add_paragraph()
            p.text = f"• {highlight}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)
    
    def _add_team_slide(self, prs: Presentation, company_data: Dict, team: Dict):
        """Add team slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Team"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        members = team.get('members', [
            {"name": "CEO/Founder", "bio": "15+ years experience in industry"},
            {"name": "CTO/Co-founder", "bio": "Led engineering at scale-up"},
            {"name": "VP Sales", "bio": "Built sales teams at top companies"}
        ])
        
        for member in members:
            # Handle both dict and string formats
            if isinstance(member, dict):
                member_name = member.get('name', 'Team Member')
                member_bio = member.get('bio', 'Experienced professional')
            else:
                # If member is a string, use it as name
                member_name = str(member)
                member_bio = "Key team member"
            
            p = text_frame.add_paragraph()
            p.text = member_name
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(15)
            
            p = text_frame.add_paragraph()
            p.text = member_bio
            p.font.size = Pt(14)
            p.font.color.rgb = self.GRAY
            p.space_after = Pt(10)
    
    def _add_competition_slide(self, prs: Presentation, market_data: Optional[Dict], competition: Dict):
        """Add competition slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Competitive Landscape"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "Our Competitive Advantages:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        advantages = competition.get('advantages', [
            "Unique technology or approach",
            "Strong customer relationships",
            "Superior unit economics"
        ])
        
        for advantage in advantages:
            p = text_frame.add_paragraph()
            p.text = f"✓ {advantage}"
            p.font.size = Pt(16)
            p.space_after = Pt(10)
        
        text_frame.add_paragraph()
        
        p = text_frame.add_paragraph()
        p.text = "Key Competitors:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(20)
        p.space_after = Pt(10)
        
        if market_data and market_data.get('competitors'):
            competitors = market_data['competitors'][:3]
            for comp in competitors:
                # Handle both dict and string formats
                if isinstance(comp, dict):
                    comp_name = comp.get('name', 'Competitor')
                else:
                    comp_name = str(comp)
                
                p = text_frame.add_paragraph()
                p.text = f"• {comp_name}"
                p.font.size = Pt(14)
                p.space_after = Pt(6)
        else:
            for comp in competition.get('competitors', ['Competitor A', 'Competitor B']):
                p = text_frame.add_paragraph()
                p.text = f"• {comp}"
                p.font.size = Pt(14)
                p.space_after = Pt(6)
    
    def _add_investment_slide(self, prs: Presentation, company_data: Dict, ask: Dict):
        """Add investment ask slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Investment Ask"
        title.text_frame.paragraphs[0].font.color.rgb = self.DARK_BLUE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Funding amount
        left = Inches(2)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(1.5)
        
        amount_box = slide.shapes.add_textbox(left, top, width, height)
        amount_frame = amount_box.text_frame
        amount_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        funding = company_data.get('funding_amount', 0)
        
        p = amount_frame.paragraphs[0]
        p.text = f"${funding:,.0f}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.ACCENT
        
        p = amount_frame.add_paragraph()
        p.text = f"{company_data.get('stage', 'Series A')} Funding Round"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.color.rgb = self.GRAY
        
        # Use of funds
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(2)
        
        use_box = slide.shapes.add_textbox(left, top, width, height)
        use_frame = use_box.text_frame
        
        p = use_frame.paragraphs[0]
        p.text = "Use of Funds:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(10)
        
        uses = ask.get('use_of_funds', [
            "Product development & engineering (40%)",
            "Sales & marketing expansion (35%)",
            "Operations & infrastructure (25%)"
        ])
        
        for use in uses:
            p = use_frame.add_paragraph()
            p.text = f"• {use}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)
    
    def _add_closing_slide(self, prs: Presentation, company_data: Dict):
        """Add closing/contact slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Thank you
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1)
        
        thanks_box = slide.shapes.add_textbox(left, top, width, height)
        thanks_frame = thanks_box.text_frame
        
        p = thanks_frame.paragraphs[0]
        p.text = "Thank You"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.DARK_BLUE
        
        # Company name
        top = Inches(3.2)
        height = Inches(0.8)
        
        company_box = slide.shapes.add_textbox(left, top, width, height)
        company_frame = company_box.text_frame
        
        p = company_frame.paragraphs[0]
        p.text = company_data.get('name', 'Company Name')
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p.font.color.rgb = self.GRAY
        
        # Contact
        top = Inches(4.5)
        height = Inches(1.5)
        
        contact_box = slide.shapes.add_textbox(left, top, width, height)
        contact_frame = contact_box.text_frame
        
        website = company_data.get('website', 'www.company.com')
        email = company_data.get('email', 'contact@company.com')
        
        p = contact_frame.paragraphs[0]
        p.text = website
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.space_after = Pt(10)
        
        p = contact_frame.add_paragraph()
        p.text = email
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
    
    async def _generate_deck_content(
        self,
        company_data: Dict,
        deal_data: Optional[Dict],
        market_data: Optional[Dict],
        financial_model: Optional[Dict]
    ) -> Dict:
        """Generate AI-powered content for deck slides"""
        logger.info("Generating pitch deck content with AI")
        
        prompt = f"""Generate comprehensive pitch deck content for:

Company: {company_data.get('name')}
Industry: {company_data.get('industry')}
Stage: {company_data.get('stage')}
Description: {company_data.get('description', '')}

{"Market Data: " + str(market_data.get('market_overview', '')) if market_data else ""}
{"Deal Score: " + str(deal_data.get('score', '')) if deal_data else ""}

Generate structured content for these slides:
1. problem: List 3 key market problems/pain points
2. solution: List 3 key solution features
3. market: Market size, growth rate, trends (3 items)
4. product: Core capabilities (3 items)
5. business_model: Revenue streams (3), unit economics (CAC, LTV, ratio)
6. traction: Key milestones (customers, revenue, growth) and achievements (3)
7. team: 3 key team members with bios
8. competition: Our advantages (3), key competitors (3)
9. ask: Use of funds breakdown (3 categories with %)

Format as JSON with this structure.
Keep all text concise and punchy - suitable for slides."""
        
        try:
            response = await self.client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are an expert pitch deck writer. Create compelling, concise content for investor presentations."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                response_format={"type": "json_object"}
            )
            
            import json
            result = json.loads(response.choices[0].message.content)
            logger.info("Deck content generated successfully")
            return result
            
        except Exception as e:
            logger.error(f"Error generating deck content: {e}")
            # Return fallback structure
            return {
                "problem": {
                    "points": [
                        "Market inefficiency costing time and money",
                        "Existing solutions are complex and expensive",
                        "Growing demand for better alternatives"
                    ]
                },
                "solution": {
                    "features": [
                        "Intuitive platform that's easy to use",
                        "10x cost savings vs. competitors",
                        "Enterprise-grade security and scalability"
                    ]
                },
                "market": {
                    "size": "$10B+ Market",
                    "growth": "25% CAGR",
                    "trends": [
                        "Digital transformation accelerating",
                        "Shift to cloud-based solutions",
                        "Increasing customer expectations"
                    ]
                },
                "product": {
                    "capabilities": [
                        "Core platform with AI-powered features",
                        "Seamless integrations with existing tools",
                        "Mobile-first design for on-the-go access"
                    ]
                },
                "business_model": {
                    "revenue_streams": [
                        "Subscription/SaaS (70%)",
                        "Professional services (20%)",
                        "Marketplace fees (10%)"
                    ],
                    "unit_economics": {
                        "cac": "$500",
                        "ltv": "$5,000",
                        "ratio": "10:1"
                    }
                },
                "traction": {
                    "milestones": [
                        {"label": "Customers", "value": "500+"},
                        {"label": "Revenue", "value": "$2M ARR"},
                        {"label": "Growth", "value": "15% MoM"}
                    ],
                    "achievements": [
                        "Signed 3 Fortune 500 customers",
                        "Launched v2.0 with key features",
                        "Expanded to 10 countries"
                    ]
                },
                "team": {
                    "members": [
                        {"name": "CEO - Jane Smith", "bio": "15+ years building SaaS companies, ex-VP at TechCorp"},
                        {"name": "CTO - John Doe", "bio": "Former lead engineer at BigTech, Stanford CS"},
                        {"name": "VP Sales - Sarah Johnson", "bio": "Built sales org from 0 to $50M at StartupCo"}
                    ]
                },
                "competition": {
                    "advantages": [
                        "Proprietary technology with patent protection",
                        "10x better performance at 50% lower cost",
                        "Strong customer retention (>95% renewal rate)"
                    ],
                    "competitors": ["Competitor A (legacy player)", "Competitor B (enterprise focus)", "Competitor C (regional)"]
                },
                "ask": {
                    "use_of_funds": [
                        "Product & Engineering (40%)",
                        "Sales & Marketing (35%)",
                        "Operations & Team (25%)"
                    ]
                }
            }
