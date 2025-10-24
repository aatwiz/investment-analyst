"""
File processing service for document extraction and parsing
Uses PyMuPDF for PDFs, python-docx for DOCX, and pandas for Excel/CSV
"""
from pathlib import Path
from typing import Dict, Any, List
import pandas as pd
import fitz  # PyMuPDF
from docx import Document

from utils.logger import setup_logger

logger = setup_logger(__name__)


class FileProcessor:
    """Handle file processing and extraction"""
    
    def __init__(self):
        self.supported_types = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_docx,
            'xlsx': self._process_excel,
            'xls': self._process_excel,
            'csv': self._process_csv,
            'txt': self._process_text,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
        }
    
    def get_file_type(self, file_path: Path) -> str:
        """Get file type from path"""
        return file_path.suffix.lstrip('.').lower()
    
    def can_process(self, file_path: Path) -> bool:
        """Check if file type is supported"""
        file_type = self.get_file_type(file_path)
        return file_type in self.supported_types
    
    def process_file(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a file and extract content
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted content and metadata
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_type = self.get_file_type(file_path)
        
        if not self.can_process(file_path):
            raise ValueError(f"Unsupported file type: {file_type}")
        
        try:
            processor = self.supported_types[file_type]
            result = processor(file_path)
            logger.info(f"Successfully processed {file_type}: {file_path.name}")
            return result
        except Exception as e:
            logger.error(f"Error processing {file_path.name}: {str(e)}")
            raise
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """
        Process PDF files using PyMuPDF
        
        Returns:
            Dict with extracted text, tables, and metadata
        """
        logger.info(f"Processing PDF: {file_path.name}")
        
        try:
            doc = fitz.open(file_path)
            
            text_content = []
            tables = []
            metadata = doc.metadata
            
            # Extract text from each page
            for page_num, page in enumerate(doc):
                # Extract text
                page_text = page.get_text()
                text_content.append(page_text)
                
                # Extract tables (basic detection)
                tables_on_page = page.find_tables()
                if tables_on_page:
                    for table in tables_on_page:
                        try:
                            tables.append({
                                "page": page_num + 1,
                                "content": table.extract()
                            })
                        except:
                            pass
            
            full_text = "\n".join(text_content)
            
            doc.close()
            
            return {
                "type": "pdf",
                "status": "success",
                "filename": file_path.name,
                "text": full_text,
                "text_length": len(full_text),
                "tables": tables,
                "table_count": len(tables),
                "page_count": len(text_content),
                "metadata": metadata,
                "summary": {
                    "total_characters": len(full_text),
                    "total_pages": len(text_content),
                    "has_tables": len(tables) > 0
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path.name}: {str(e)}")
            return {
                "type": "pdf",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """
        Process DOCX files using python-docx
        
        Returns:
            Dict with extracted text and metadata
        """
        logger.info(f"Processing DOCX: {file_path.name}")
        
        try:
            doc = Document(file_path)
            
            # Extract text from paragraphs
            text_content = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            full_text = "\n".join(text_content)
            
            return {
                "type": "docx",
                "status": "success",
                "filename": file_path.name,
                "text": full_text,
                "text_length": len(full_text),
                "tables": tables,
                "table_count": len(tables),
                "paragraph_count": len(text_content),
                "summary": {
                    "total_characters": len(full_text),
                    "total_paragraphs": len(text_content),
                    "has_tables": len(tables) > 0
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path.name}: {str(e)}")
            return {
                "type": "docx",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """
        Process Excel files using pandas
        
        Returns:
            Dict with dataframes and metadata
        """
        logger.info(f"Processing Excel: {file_path.name}")
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            sheets_data = {}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert to dict with summary
                sheets_data[sheet_name] = {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": df.columns.tolist(),
                    "data_preview": df.head(10).to_dict(orient='records'),  # First 10 rows
                    "data_types": df.dtypes.astype(str).to_dict(),
                    "summary_stats": df.describe().to_dict() if not df.empty else {}
                }
            
            return {
                "type": "excel",
                "status": "success",
                "filename": file_path.name,
                "sheet_count": len(sheets_data),
                "sheet_names": list(sheets_data.keys()),
                "sheets": sheets_data,
                "summary": {
                    "total_sheets": len(sheets_data),
                    "total_rows": sum(s["rows"] for s in sheets_data.values()),
                    "total_columns": sum(s["columns"] for s in sheets_data.values())
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing Excel {file_path.name}: {str(e)}")
            return {
                "type": "excel",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """
        Process CSV files using pandas
        
        Returns:
            Dict with dataframe and metadata
        """
        logger.info(f"Processing CSV: {file_path.name}")
        
        try:
            df = pd.read_csv(file_path)
            
            return {
                "type": "csv",
                "status": "success",
                "filename": file_path.name,
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": df.columns.tolist(),
                "data_preview": df.head(10).to_dict(orient='records'),
                "data_types": df.dtypes.astype(str).to_dict(),
                "summary_stats": df.describe().to_dict() if not df.empty else {},
                "summary": {
                    "total_rows": len(df),
                    "total_columns": len(df.columns)
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing CSV {file_path.name}: {str(e)}")
            return {
                "type": "csv",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Process text files"""
        logger.info(f"Processing TXT: {file_path.name}")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                "type": "text",
                "status": "success",
                "filename": file_path.name,
                "text": content,
                "text_length": len(content),
                "line_count": content.count('\n') + 1,
                "word_count": len(content.split()),
                "summary": {
                    "total_characters": len(content),
                    "total_lines": content.count('\n') + 1,
                    "total_words": len(content.split())
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing TXT {file_path.name}: {str(e)}")
            return {
                "type": "text",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process PPTX files using python-pptx"""
        logger.info(f"Processing PPTX: {file_path.name}")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            full_text = "\n".join(text_content)
            
            return {
                "type": "pptx",
                "status": "success",
                "filename": file_path.name,
                "text": full_text,
                "text_length": len(full_text),
                "slide_count": len(prs.slides),
                "summary": {
                    "total_characters": len(full_text),
                    "total_slides": len(prs.slides)
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path.name}: {str(e)}")
            return {
                "type": "pptx",
                "status": "error",
                "filename": file_path.name,
                "error": str(e)
            }
    
    def _process_ppt(self, file_path: Path) -> Dict[str, Any]:
        """Process PPT files"""
        logger.info(f"Processing PPT: {file_path.name}")
        return self._process_pptx(file_path)  # Use same method
